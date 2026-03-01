#!/usr/bin/env python3
"""
Fetches the last 2 weeks of GitHub changelog updates from the RSS feed
and generates a PowerPoint presentation.

Japanese translation is handled by the Copilot agent (via the changelog-ppt-creator
skill), not by this script. The script leaves title_ja/bullets_ja empty for the
agent to fill in.

Usage:
    python3 scripts/fetch-changelog.py [--days 14] [--output-dir changelog]

Requirements:
    pip3 install python-pptx requests beautifulsoup4 Pillow
"""

from __future__ import annotations

import argparse
import os
import re
import sys
import xml.etree.ElementTree as ET
from datetime import datetime, timedelta, timezone
from io import BytesIO
from pathlib import Path
from typing import Optional

import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Pt, Emu

FEED_URL = "https://github.blog/changelog/feed/"
REPO_ROOT = Path(__file__).resolve().parent.parent
TEMPLATE_PATH = REPO_ROOT / ".github" / "skills" / "changelog-ppt-creator" / "changelog-template.pptx"

NAMESPACES = {
    "content": "http://purl.org/rss/1.0/modules/content/",
    "dc": "http://purl.org/dc/elements/1.1/",
}


def parse_date(date_str: str) -> datetime:
    """Parse RSS pubDate format."""
    # e.g. "Thu, 19 Feb 2026 22:39:37 +0000"
    from email.utils import parsedate_to_datetime
    return parsedate_to_datetime(date_str)


def fetch_feed(days: int) -> list[dict]:
    """Fetch RSS feed (paginated) and return entries from the last N days."""
    cutoff = datetime.now(timezone.utc) - timedelta(days=days)
    entries = []

    for page in range(1, 20):  # Safety limit
        url = FEED_URL if page == 1 else f"{FEED_URL}?paged={page}"
        resp = requests.get(url, timeout=30)
        if resp.status_code != 200:
            break

        root = ET.fromstring(resp.text)
        items = root.findall(".//item")
        if not items:
            break

        page_has_old = False
        for item in items:
            pub_date = parse_date(item.find("pubDate").text)
            if pub_date < cutoff:
                page_has_old = True
                continue

            title = item.find("title").text
            link = item.find("link").text
            content_el = item.find("content:encoded", NAMESPACES)
            content_html = content_el.text if content_el is not None else ""

            categories = []
            for cat in item.findall("category"):
                categories.append({
                    "domain": cat.get("domain", ""),
                    "name": cat.text,
                })

            label = next((c["name"] for c in categories if c["domain"] == "changelog-label"), "other")
            change_type = next((c["name"] for c in categories if c["domain"] == "changelog-type"), "")

            entries.append({
                "title": title,
                "link": link,
                "date": pub_date,
                "content_html": content_html,
                "label": label,
                "type": change_type,
            })

        # Stop if we've passed the cutoff date
        if page_has_old:
            break

    entries.sort(key=lambda e: e["date"], reverse=True)
    return entries


def extract_summary_and_images(html: str) -> "tuple[list[str], list[str]]":
    """Extract bullet-point summary and image URLs from HTML content."""
    soup = BeautifulSoup(html, "html.parser")

    # Images from RSS content are rare; they'll be supplemented by scrape_page_images
    images = []
    for img in soup.find_all("img"):
        src = img.get("src", "")
        if src and "favicon" not in src and "avatar" not in src and not src.startswith("data:"):
            images.append(src)

    # Extract summary bullet points
    bullets = []

    # Get all list items first
    for li in soup.find_all("li"):
        text = li.get_text(strip=True)
        if text and len(text) > 10:
            bullets.append(text)

    # If no list items, get first few paragraphs
    if not bullets:
        for p in soup.find_all("p"):
            text = p.get_text(strip=True)
            if text and len(text) > 20:
                # Skip "The post ... appeared first on ..." boilerplate
                if "appeared first on" in text or "Join the discussion" in text:
                    continue
                bullets.append(text)
                if len(bullets) >= 4:
                    break

    # Limit bullets and truncate to fit slides (Japanese ~50 chars/line at 20pt)
    bullets = [b[:120] + "…" if len(b) > 120 else b for b in bullets[:5]]

    # Extract full plain text for slide notes
    plain_text = soup.get_text(separator="\n", strip=True)
    # Remove boilerplate
    plain_lines = []
    for line in plain_text.split("\n"):
        line = line.strip()
        if not line or "appeared first on" in line or "Join the discussion" in line:
            continue
        plain_lines.append(line)

    return bullets, images, "\n".join(plain_lines)


def scrape_page_images(url: str) -> list:
    """Scrape the actual changelog page for images and og:image."""
    images = []
    try:
        resp = requests.get(url, timeout=15, headers={"User-Agent": "Mozilla/5.0"})
        resp.raise_for_status()
        soup = BeautifulSoup(resp.text, "html.parser")

        # Look for content images in article/main
        article = soup.find("article") or soup.find("main") or soup
        for img in article.find_all("img"):
            src = img.get("src", "")
            if src and "favicon" not in src and "avatar" not in src and "logo" not in src and not src.startswith("data:"):
                images.append(src)

        # Fallback to og:image
        if not images:
            og = soup.find("meta", property="og:image")
            if og and og.get("content"):
                images.append(og["content"])
    except Exception:
        pass
    return images


def download_image(url: str) -> Optional[bytes]:
    """Download an image and return its bytes. Validates it's a real image."""
    try:
        # Clean URL (remove resize params that might cause issues)
        clean_url = url.split("?")[0] if "?" in url else url
        resp = requests.get(url, timeout=15, headers={"User-Agent": "Mozilla/5.0"})
        resp.raise_for_status()
        content_type = resp.headers.get("content-type", "")
        if "image" in content_type:
            # Validate it's a real image PIL can read
            from PIL import Image as PILImage
            PILImage.open(BytesIO(resp.content)).verify()
            return resp.content
    except Exception:
        pass
    return None


def _clean_label(label: str) -> str:
    """Clean label text: decode HTML entities and replace & with 'and'."""
    import html
    label = html.unescape(label)
    label = label.replace("&", "and")
    return label


# Consolidate granular RSS labels into broader presentation categories
LABEL_GROUP_MAP = {
    "actions": "Actions",
    "copilot": "Copilot",
    "collaboration tools": "Collaboration",
    "projects and issues": "Collaboration",
    "projects & issues": "Collaboration",
    "client apps": "Copilot",
    "account management": "Administration",
    "enterprise management tools": "Administration",
    "platform governance": "Administration",
    "application security": "Security",
    "supply chain security": "Security",
    "code security": "Security",
    "secret scanning": "Security",
    "dependabot": "Security",
    "code scanning": "Security",
}

# Preferred display order
LABEL_ORDER = ["Actions", "Copilot", "Collaboration", "Security", "Administration"]


def _map_label(label: str) -> str:
    """Map a raw RSS label to a consolidated presentation category."""
    cleaned = _clean_label(label).lower()
    return LABEL_GROUP_MAP.get(cleaned, _clean_label(label).title())


def group_by_label(entries: list[dict]) -> dict[str, list[dict]]:
    """Group changelog entries by consolidated label."""
    groups = {}
    for entry in entries:
        label = _map_label(entry["label"])
        groups.setdefault(label, []).append(entry)
    # Sort by preferred order, then alphabetically for unknown categories
    def sort_key(label):
        if label in LABEL_ORDER:
            return (0, LABEL_ORDER.index(label))
        return (1, label)
    return dict(sorted(groups.items(), key=lambda x: sort_key(x[0])))


def generate_markdown(entries: list[dict], lang: str = "en") -> str:
    """Generate a markdown summary of changelog entries."""
    today = datetime.now().strftime("%Y-%m-%d")
    if lang == "ja":
        lines = [f"# GitHub Changelog アップデート\n", f"**生成日**: {today}\n"]
    else:
        lines = [f"# GitHub Changelog Updates\n", f"**Generated**: {today}\n"]

    for entry in entries:
        lines.append(f"## {entry['title_ja'] if lang == 'ja' else entry['title']}\n")
        lines.append(f"- **Date**: {entry['date'].strftime('%Y-%m-%d')}")
        lines.append(f"- **Type**: {entry['type']}")
        lines.append(f"- **Category**: {entry['label']}")
        lines.append(f"- **Link**: {entry['link']}\n")

        bullets = entry.get("bullets_ja" if lang == "ja" else "bullets", [])
        for b in bullets:
            lines.append(f"- {b}")
        lines.append("")

    return "\n".join(lines)


def _add_run(paragraph, text, bold=True, size=None, color_rgb=None):
    """Add a formatted run to a paragraph, matching template styling."""
    run = paragraph.add_run()
    run.text = text
    run.font.bold = bold
    if size:
        run.font.size = size
    if color_rgb:
        run.font.color.rgb = color_rgb


def build_pptx(entries: list[dict], output_path: Path, lang: str = "ja"):
    """Build a PPTX by cloning template slides and replacing text/images.

    Args:
        lang: "ja" for Japanese, "en" for English

    Template slide map:
      0 = Title slide
      1 = TOC
      2 = Section header
      3 = Content slide (summary + date/category tiles) — IMPROVEMENT
      4 = Image slide (full image + date/category tiles)
      5 = Content slide — RETIRED category tile
      6 = Content slide — RELEASE category tile
      7 = Section header (alternate)
    """
    import copy
    from lxml import etree
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    prs = Presentation(str(TEMPLATE_PATH))

    NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    # Store all template slide XMLs before modifying
    tpl_xmls = [copy.deepcopy(s._element) for s in prs.slides]

    # Store all non-layout template relationships for cloning
    # Each entry: {rId: (reltype, target_part_or_url, is_external)}
    tpl_extra_rels = {}
    for idx, slide in enumerate(prs.slides):
        extra = {}
        for rel in slide.part.rels.values():
            if "slideLayout" not in rel.reltype and "notesSlide" not in rel.reltype:
                if rel.is_external:
                    extra[rel.rId] = (rel.reltype, rel.target_ref, True)
                else:
                    extra[rel.rId] = (rel.reltype, rel.target_part, False)
        if extra:
            tpl_extra_rels[idx] = extra

    # Store template image from slide 4 for reuse
    tpl_pic_xml = None
    spTree = tpl_xmls[4].find(f".//{{{NS_P}}}spTree")
    for el in spTree:
        if el.tag == f"{{{NS_P}}}pic":
            tpl_pic_xml = copy.deepcopy(el)
            break

    # Category tile colors: pick template slide based on type
    # Slide 3 = IMPROVEMENT (#092F69), Slide 5 = RETIRED (#67070C), Slide 6 = RELEASE (#053A16)
    CATEGORY_TILE_COLORS = {
        "Improvement": "092F69",
        "Release": "053A16",
        "Retired": "67070C",
    }
    # Map changelog type to the correct template slide index for category tile
    CATEGORY_TPL_SLIDE = {
        "Improvement": 3,
        "Release": 6,
        "Retired": 5,
    }

    def _find_shapes(slide_xml):
        sp_tree = slide_xml.find(f".//{{{NS_P}}}spTree")
        return list(sp_tree) if sp_tree is not None else []

    def _get_shape_name(shape_xml):
        cNvPr = shape_xml.find(f".//{{{NS_P}}}cNvPr")
        if cNvPr is None:
            cNvPr = shape_xml.find(f".//{{{NS_A}}}cNvPr")
        return cNvPr.get("name", "") if cNvPr is not None else ""

    def _find_paragraphs(shape_xml):
        txBody = shape_xml.find(f".//{{{NS_P}}}txBody")
        if txBody is None:
            txBody = shape_xml.find(f".//{{{NS_A}}}txBody")
        return txBody.findall(f"{{{NS_A}}}p") if txBody is not None else []

    def _clone_slide(prs, tpl_xml, tpl_idx):
        layout_idx = 0 if tpl_idx == 0 else 1
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        new_xml = copy.deepcopy(tpl_xml)
        el = slide._element
        for child in list(el):
            el.remove(child)
        for key, val in new_xml.attrib.items():
            el.set(key, val)
        for child in list(new_xml):
            el.append(child)
        # Copy non-layout relationships from template slide (images, hyperlinks)
        if tpl_idx in tpl_extra_rels:
            from pptx.opc.package import _Relationship
            for rId, (reltype, target, is_external) in tpl_extra_rels[tpl_idx].items():
                if is_external:
                    slide.part.rels._rels[rId] = _Relationship(
                        slide.part.rels._base_uri,
                        rId,
                        reltype,
                        target_mode="External",
                        target=target,
                    )
                else:
                    slide.part.rels._rels[rId] = _Relationship(
                        slide.part.rels._base_uri,
                        rId,
                        reltype,
                        target_mode="Internal",
                        target=target,
                    )
        return slide

    def _set_title_text(slide_xml, text):
        for shape in _find_shapes(slide_xml):
            if "Title" in _get_shape_name(shape):
                paras = _find_paragraphs(shape)
                if paras:
                    runs = paras[0].findall(f"{{{NS_A}}}r")
                    if runs:
                        runs[0].find(f"{{{NS_A}}}t").text = text
                        for r in runs[1:]:
                            paras[0].remove(r)
                return

    def _set_subtitle_texts(slide_xml, texts):
        for shape in _find_shapes(slide_xml):
            if "Subtitle" in _get_shape_name(shape):
                paras = _find_paragraphs(shape)
                if paras:
                    runs = paras[0].findall(f"{{{NS_A}}}r")
                    for i, run in enumerate(runs):
                        t_elem = run.find(f"{{{NS_A}}}t")
                        if i < len(texts):
                            t_elem.text = texts[i]
                        else:
                            paras[0].remove(run)
                return

    def _set_content_bullets(slide_xml, bullets):
        """Replace content placeholder with bullet list, using template paragraph formatting."""
        for shape in _find_shapes(slide_xml):
            if "Content" in _get_shape_name(shape):
                txBody = shape.find(f".//{{{NS_P}}}txBody")
                if txBody is None:
                    continue
                paras = txBody.findall(f"{{{NS_A}}}p")
                if not paras:
                    continue
                # Use first content paragraph as stencil
                bullet_tpl = copy.deepcopy(paras[0])
                # Also keep the empty separator if exists
                sep_tpl = copy.deepcopy(paras[1]) if len(paras) > 1 else None
                # Remove all paragraphs
                for p in paras:
                    txBody.remove(p)
                # Add new bullet paragraphs
                for i, bullet_text in enumerate(bullets):
                    if i > 0 and sep_tpl is not None:
                        txBody.append(copy.deepcopy(sep_tpl))
                    p = copy.deepcopy(bullet_tpl)
                    runs = p.findall(f"{{{NS_A}}}r")
                    if runs:
                        runs[0].find(f"{{{NS_A}}}t").text = bullet_text
                        for r in runs[1:]:
                            p.remove(r)
                    txBody.append(p)
                return

    def _set_date_tile(slide_xml, date_str):
        """Set the date tile text (first Rounded Rectangle)."""
        for shape in _find_shapes(slide_xml):
            name = _get_shape_name(shape)
            if "Rounded" in name:
                paras = _find_paragraphs(shape)
                if paras:
                    runs = paras[0].findall(f"{{{NS_A}}}r")
                    if runs and len(runs) >= 3:
                        # Template has "FEBRUARY" " " "23, 1995"
                        month = date_str.strftime("%B").upper()
                        day_year = date_str.strftime("%d, %Y")
                        runs[0].find(f"{{{NS_A}}}t").text = month
                        runs[1].find(f"{{{NS_A}}}t").text = " "
                        runs[2].find(f"{{{NS_A}}}t").text = day_year
                        for r in runs[3:]:
                            paras[0].remove(r)
                    elif runs:
                        runs[0].find(f"{{{NS_A}}}t").text = date_str.strftime("%B %d, %Y").upper()
                return  # Only first Rounded Rectangle is the date

    def _set_category_tile(slide_xml, category_type, tpl_xmls):
        """Replace the category tile with the correct color variant from templates."""
        tpl_idx = CATEGORY_TPL_SLIDE.get(category_type, 3)
        # Find the category tile in the template (second Rounded Rectangle)
        src_tile = None
        rounded_count = 0
        for shape in _find_shapes(tpl_xmls[tpl_idx]):
            name = _get_shape_name(shape)
            if "Rounded" in name:
                rounded_count += 1
                if rounded_count == 2:
                    src_tile = copy.deepcopy(shape)
                    break

        if src_tile is None:
            return

        # Find and replace the category tile in the target slide
        sp_tree = slide_xml.find(f".//{{{NS_P}}}spTree")
        rounded_count = 0
        for shape in list(_find_shapes(slide_xml)):
            name = _get_shape_name(shape)
            if "Rounded" in name:
                rounded_count += 1
                if rounded_count == 2:
                    # Replace with correct color variant
                    sp_tree.remove(shape)
                    # Set text on the source tile
                    paras = src_tile.findall(f".//{{{NS_A}}}p")
                    if paras:
                        runs = paras[0].findall(f"{{{NS_A}}}r")
                        if runs:
                            runs[0].find(f"{{{NS_A}}}t").text = category_type.upper()
                            for r in runs[1:]:
                                paras[0].remove(r)
                    sp_tree.append(src_tile)
                    return

    def _remove_image(slide_xml):
        sp_tree = slide_xml.find(f".//{{{NS_P}}}spTree")
        if sp_tree is not None:
            for pic in sp_tree.findall(f"{{{NS_P}}}pic"):
                sp_tree.remove(pic)

    def _add_image_from_template(slide, tpl_pic_xml, image_bytes):
        from pptx.parts.image import ImagePart, Image as PptxImage
        try:
            pic_xml = copy.deepcopy(tpl_pic_xml)
            image = PptxImage.from_blob(image_bytes)
            image_part = ImagePart.new(prs.part.package, image)
            rId = slide.part.relate_to(image_part, RT.IMAGE)
            blip = pic_xml.find(f".//{{{NS_A}}}blip")
            if blip is not None:
                blip.set(f"{{{NS_R}}}embed", rId)
            spTree = slide._element.find(f".//{{{NS_P}}}spTree")
            spTree.append(pic_xml)
        except Exception:
            pass

    def _set_slide_notes(slide, text):
        """Add text to the slide's notes section."""
        from pptx.util import Pt as _Pt
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.clear()
        tf.text = text

    # --- Remove all existing slides ---
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(f"{{{NS_R}}}id")
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    # --- Slide 0: Title ---
    title_xml = copy.deepcopy(tpl_xmls[0])
    dates = [e["date"] for e in entries]
    if dates:
        if lang == "ja":
            _set_subtitle_texts(title_xml, [
                min(dates).strftime("%-m月%-d日"), "〜 ",
                max(dates).strftime("%-m月%-d日"), " ", "のアップデート"
            ])
        else:
            _set_subtitle_texts(title_xml, [
                min(dates).strftime("%B"), " ",
                min(dates).strftime("%d"), " - ",
                max(dates).strftime("%B %d, %Y")
            ])
    _clone_slide(prs, title_xml, 0)

    # --- Slide 1: TOC ---
    groups = group_by_label(entries)
    toc_xml = copy.deepcopy(tpl_xmls[1])
    for shape in _find_shapes(toc_xml):
        if "Content" in _get_shape_name(shape):
            txBody = shape.find(f".//{{{NS_P}}}txBody")
            if txBody is None:
                continue
            paras = txBody.findall(f"{{{NS_A}}}p")
            toc_tpl = None
            empty_tpl = None
            for p in paras:
                runs = p.findall(f"{{{NS_A}}}r")
                if runs and toc_tpl is None:
                    toc_tpl = copy.deepcopy(p)
                elif not runs and empty_tpl is None:
                    empty_tpl = copy.deepcopy(p)
            for p in paras:
                txBody.remove(p)
            for i, (label, items) in enumerate(groups.items(), 1):
                if i > 1 and empty_tpl is not None:
                    txBody.append(copy.deepcopy(empty_tpl))
                p = copy.deepcopy(toc_tpl)
                runs = p.findall(f"{{{NS_A}}}r")
                if runs:
                    count_text = f"{len(items)}件" if lang == "ja" else str(len(items))
                    runs[0].find(f"{{{NS_A}}}t").text = f"{i} {label.title()} ("
                    if len(runs) > 1:
                        runs[1].find(f"{{{NS_A}}}t").text = count_text
                    if len(runs) > 2:
                        runs[2].find(f"{{{NS_A}}}t").text = ")"
                    for r in runs[3:]:
                        p.remove(r)
                txBody.append(p)
            break
    _clone_slide(prs, toc_xml, 1)

    # --- Per-category slides ---
    # Track sections: list of (name, first_slide_index)
    sections = [("Introduction", 0)]  # Title + TOC

    for label, items in groups.items():
        section_start = len(prs.slides)

        # Section header
        section_xml = copy.deepcopy(tpl_xmls[2])
        count_suffix = f"{len(items)}件" if lang == "ja" else str(len(items))
        _set_title_text(section_xml, f"{label.title()} ({count_suffix})")
        _clone_slide(prs, section_xml, 2)

        for entry in items:
            cat_type = entry.get("type", "Improvement")
            entry_date = entry["date"]

            # Pick title and bullets based on language
            title = entry.get("title_ja", entry["title"]) if lang == "ja" else entry["title"]
            bullets = entry.get("bullets_ja", []) if lang == "ja" else entry.get("bullets", [])

            # --- Summary slide (clone from slide 3) ---
            content_xml = copy.deepcopy(tpl_xmls[3])
            _set_title_text(content_xml, title)
            _set_content_bullets(content_xml, bullets)
            _set_date_tile(content_xml, entry_date)
            _set_category_tile(content_xml, cat_type, tpl_xmls)
            summary_slide = _clone_slide(prs, content_xml, 3)

            # Add full English content to slide notes
            notes_text = f"{entry['title']}\n{entry['link']}\n\n{entry.get('plain_text_en', '')}"
            _set_slide_notes(summary_slide, notes_text)

            # --- Image slide (clone from slide 4) if image available ---
            img_data_list = entry.get("image_data", [])
            if img_data_list and tpl_pic_xml is not None:
                img_xml = copy.deepcopy(tpl_xmls[4])
                _set_title_text(img_xml, title)
                _set_date_tile(img_xml, entry_date)
                _set_category_tile(img_xml, cat_type, tpl_xmls)
                _remove_image(img_xml)
                slide = _clone_slide(prs, img_xml, 4)
                _add_image_from_template(slide, tpl_pic_xml, img_data_list[0])

        sections.append((f"{label.title()} ({len(items)})", section_start))

    # --- Add real PowerPoint sections to presentation.xml ---
    # Sections live inside <p:extLst><p:ext uri="{521415D9-...}"><p14:sectionLst>
    import uuid
    prs_element = prs.part._element
    NS_P14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
    SECTION_EXT_URI = "{521415D9-36F7-43E2-AB2F-B90AF26B5E84}"

    # Get slide IDs from sldIdLst
    sld_id_lst = prs_element.find(f"{{{NS_P}}}sldIdLst")
    slide_ids = [el.get("id") for el in sld_id_lst]

    # Find the existing p14:sectionLst inside extLst
    section_lst = None
    ext_lst = prs_element.find(f"{{{NS_P}}}extLst")
    if ext_lst is not None:
        for ext in ext_lst.findall(f"{{{NS_P}}}ext"):
            if ext.get("uri") == SECTION_EXT_URI:
                section_lst = ext.find(f"{{{NS_P14}}}sectionLst")
                if section_lst is None:
                    # Create it inside this ext
                    section_lst = etree.SubElement(ext, f"{{{NS_P14}}}sectionLst")
                else:
                    # Clear existing sections
                    for child in list(section_lst):
                        section_lst.remove(child)
                break
        else:
            # No ext with this URI — create one
            ext_el = etree.SubElement(ext_lst, f"{{{NS_P}}}ext")
            ext_el.set("uri", SECTION_EXT_URI)
            section_lst = etree.SubElement(ext_el, f"{{{NS_P14}}}sectionLst")

    if section_lst is None:
        # No extLst at all — create the full structure
        ext_lst = etree.SubElement(prs_element, f"{{{NS_P}}}extLst")
        ext_el = etree.SubElement(ext_lst, f"{{{NS_P}}}ext")
        ext_el.set("uri", SECTION_EXT_URI)
        section_lst = etree.SubElement(ext_el, f"{{{NS_P14}}}sectionLst")

    # Build section entries
    for sec_name, sec_start in sections:
        sec_el = etree.SubElement(section_lst, f"{{{NS_P14}}}section")
        sec_el.set("name", sec_name)
        sec_el.set("id", "{" + str(uuid.uuid4()).upper() + "}")
        sld_id_lst_el = etree.SubElement(sec_el, f"{{{NS_P14}}}sldIdLst")
        # Find end of this section
        sec_idx = sections.index((sec_name, sec_start))
        if sec_idx + 1 < len(sections):
            sec_end = sections[sec_idx + 1][1]
        else:
            sec_end = len(slide_ids)
        for sid in slide_ids[sec_start:sec_end]:
            sld_id_el = etree.SubElement(sld_id_lst_el, f"{{{NS_P14}}}sldId")
            sld_id_el.set("id", sid)

    prs.save(str(output_path))


def main():
    parser = argparse.ArgumentParser(description="Fetch GitHub changelog and create PPT")
    parser.add_argument("--days", type=int, default=14, help="Number of days to look back (default: 14)")
    parser.add_argument("--output-dir", type=str, default=None, help="Output directory")
    args = parser.parse_args()

    output_dir = Path(args.output_dir) if args.output_dir else REPO_ROOT / "changelog"
    en_dir = output_dir / "english"
    ja_dir = output_dir / "japanese"
    ppt_en_dir = output_dir / "ppt-english"
    ppt_ja_dir = output_dir / "ppt-japanese"

    for d in [en_dir, ja_dir, ppt_en_dir, ppt_ja_dir]:
        d.mkdir(parents=True, exist_ok=True)

    # Fetch entries
    print(f"Fetching changelog entries from the last {args.days} days...")
    entries = fetch_feed(args.days)
    print(f"Found {len(entries)} entries.")

    if not entries:
        print("No entries found. Exiting.")
        sys.exit(0)

    # Extract summaries and images
    print("Extracting summaries and images...")
    for entry in entries:
        bullets, image_urls, plain_text_en = extract_summary_and_images(entry["content_html"])
        entry["bullets"] = bullets
        entry["plain_text_en"] = plain_text_en

        # Scrape actual page for images (RSS content rarely has them)
        if not image_urls:
            print(f"  Scraping images from: {entry['link'][:60]}...")
            image_urls = scrape_page_images(entry["link"])
        entry["image_urls"] = image_urls

        # Download first valid image (try up to 3 URLs)
        entry["image_data"] = []
        for url in image_urls[:3]:
            img_data = download_image(url)
            if img_data:
                entry["image_data"].append(img_data)
                break

    # Leave Japanese fields empty — agent fills these in
    print("Japanese fields left empty (agent handles translation).")
    for entry in entries:
        entry["title_ja"] = entry["title"]  # Fallback to English
        entry["bullets_ja"] = entry.get("bullets", [])

    # Date range for filenames (oldest to newest entry)
    dates = [e["date"] for e in entries]
    date_from = min(dates).strftime("%Y-%m-%d")
    date_to = max(dates).strftime("%Y-%m-%d")
    file_stem = f"{date_from}_to_{date_to}"

    # Generate English markdown
    en_md = generate_markdown(entries, "en")
    en_path = en_dir / f"{file_stem}-changelog.md"
    en_path.write_text(en_md, encoding="utf-8")
    print(f"Saved English markdown: {en_path}")

    # Generate Japanese markdown
    ja_md = generate_markdown(entries, "ja")
    ja_path = ja_dir / f"{file_stem}-changelog.md"
    ja_path.write_text(ja_md, encoding="utf-8")
    print(f"Saved Japanese markdown: {ja_path}")

    # Generate English PowerPoint
    ppt_en_path = ppt_en_dir / f"{file_stem}-changelog.pptx"
    print("Building English PowerPoint presentation...")
    build_pptx(entries, ppt_en_path, lang="en")
    print(f"Saved English PowerPoint: {ppt_en_path}")

    # Generate Japanese PowerPoint
    ppt_ja_path = ppt_ja_dir / f"{file_stem}-changelog.pptx"
    print("Building Japanese PowerPoint presentation...")
    build_pptx(entries, ppt_ja_path, lang="ja")
    print(f"Saved Japanese PowerPoint: {ppt_ja_path}")

    print(f"\nDone! Generated {len(entries)} changelog entries.")
    print(f"  English markdown: {en_path}")
    print(f"  Japanese markdown: {ja_path}")
    print(f"  English PowerPoint: {ppt_en_path}")
    print(f"  Japanese PowerPoint: {ppt_ja_path}")


if __name__ == "__main__":
    main()
