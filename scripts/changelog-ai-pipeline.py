#!/usr/bin/env python3
"""
All-in-one changelog pipeline: fetch RSS → AI translate → build markdown + PPTX.

Uses GitHub Models API for high-quality Japanese translation.
Designed to run in a standard GitHub Actions workflow (no firewall restrictions).

Usage:
    python3 scripts/changelog-ai-pipeline.py [--days 14] [--output-dir 10-news/11-changelog]

Dependencies:
    pip3 install python-pptx requests beautifulsoup4 Pillow lxml

Environment:
    GITHUB_TOKEN — required for GitHub Models API translation
"""

from __future__ import annotations

import argparse
import copy
import hashlib
import html as html_module
import json
import os
import sys
import uuid
import xml.etree.ElementTree as ET
from datetime import datetime, timedelta, timezone
from io import BytesIO
from pathlib import Path
from typing import Optional

import requests
from bs4 import BeautifulSoup
from lxml import etree
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import _Relationship

FEED_URL = "https://github.blog/changelog/feed/"
REPO_ROOT = Path(__file__).resolve().parent.parent
TEMPLATE_PATH = REPO_ROOT / ".github" / "skills" / "changelog-ppt-creator" / "changelog-template.pptx"
MODELS_API_URL = "https://models.github.ai/inference/chat/completions"
MODELS_MODEL = "openai/gpt-4o"

# XML namespaces
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
NS_P14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
SECTION_EXT_URI = "{521415D9-36F7-43E2-AB2F-B90AF26B5E84}"

NAMESPACES = {"content": "http://purl.org/rss/1.0/modules/content/"}

# --- Label / category mapping ---

LABEL_GROUP_MAP = {
    "actions": "Actions",
    "copilot": "Copilot",
    "collaboration tools": "Collaboration",
    "projects and issues": "Collaboration",
    "projects & issues": "Collaboration",
    "client apps": "Copilot",
    "account management": "Administration",
    "enterprise management tools": "Administration",
    "platform governance": "Administration",
    "application security": "Security",
    "supply chain security": "Security",
    "code security": "Security",
    "secret scanning": "Security",
    "dependabot": "Security",
    "code scanning": "Security",
}

LABEL_ORDER = ["Actions", "Copilot", "Collaboration", "Security", "Administration"]

CATEGORY_TPL_SLIDE = {"Improvement": 3, "Release": 6, "Retired": 5}


def _clean_label(label: str) -> str:
    label = html_module.unescape(label)
    label = label.replace("&", "and")
    return label


def _map_label(label: str) -> str:
    cleaned = _clean_label(label).lower()
    return LABEL_GROUP_MAP.get(cleaned, _clean_label(label).title())


def group_by_label(entries: list[dict]) -> dict[str, list[dict]]:
    groups: dict[str, list[dict]] = {}
    for entry in entries:
        label = _map_label(entry["label"])
        groups.setdefault(label, []).append(entry)

    def sort_key(label):
        if label in LABEL_ORDER:
            return (0, LABEL_ORDER.index(label))
        return (1, label)

    return dict(sorted(groups.items(), key=lambda x: sort_key(x[0])))


# ============================================================================
# Phase 1: RSS Fetching
# ============================================================================


def parse_date(date_str: str) -> datetime:
    from email.utils import parsedate_to_datetime
    return parsedate_to_datetime(date_str)


def fetch_feed(days: int) -> list[dict]:
    """Fetch RSS feed (paginated) and return entries from the last N days."""
    cutoff = datetime.now(timezone.utc) - timedelta(days=days)
    entries: list[dict] = []

    for page in range(1, 20):
        url = FEED_URL if page == 1 else f"{FEED_URL}?paged={page}"
        resp = requests.get(url, timeout=30)
        if resp.status_code != 200:
            break

        root = ET.fromstring(resp.text)
        items = root.findall(".//item")
        if not items:
            break

        page_has_old = False
        for item in items:
            pub_date = parse_date(item.find("pubDate").text)
            if pub_date < cutoff:
                page_has_old = True
                continue

            title = item.find("title").text
            link = item.find("link").text
            content_el = item.find("content:encoded", NAMESPACES)
            content_html = content_el.text if content_el is not None else ""

            categories = []
            for cat in item.findall("category"):
                categories.append({"domain": cat.get("domain", ""), "name": cat.text})

            label = next((c["name"] for c in categories if c["domain"] == "changelog-label"), "other")
            change_type = next((c["name"] for c in categories if c["domain"] == "changelog-type"), "")

            entries.append({
                "title": title,
                "link": link,
                "date": pub_date,
                "content_html": content_html,
                "label": label,
                "type": change_type,
            })

        if page_has_old:
            break

    entries.sort(key=lambda e: e["date"], reverse=True)
    return entries


def extract_summary_and_images(html: str) -> tuple[list[str], list[str], str]:
    """Extract bullet-point summary, image URLs, and plain text from HTML."""
    soup = BeautifulSoup(html, "html.parser")

    images = []
    for img in soup.find_all("img"):
        src = img.get("src", "")
        if src and "favicon" not in src and "avatar" not in src and not src.startswith("data:"):
            images.append(src)

    bullets = []
    for li in soup.find_all("li"):
        text = li.get_text(strip=True)
        if text and len(text) > 10:
            bullets.append(text)

    if not bullets:
        for p in soup.find_all("p"):
            text = p.get_text(strip=True)
            if text and len(text) > 20:
                if "appeared first on" in text or "Join the discussion" in text:
                    continue
                bullets.append(text)
                if len(bullets) >= 4:
                    break

    bullets = [b[:120] + "…" if len(b) > 120 else b for b in bullets[:5]]

    plain_text = soup.get_text(separator="\n", strip=True)
    plain_lines = [
        line.strip() for line in plain_text.split("\n")
        if line.strip() and "appeared first on" not in line and "Join the discussion" not in line
    ]

    return bullets, images, "\n".join(plain_lines)


def scrape_page_images(url: str) -> list[str]:
    """Scrape the actual changelog page for images and og:image."""
    images = []
    try:
        resp = requests.get(url, timeout=15, headers={"User-Agent": "Mozilla/5.0"})
        resp.raise_for_status()
        soup = BeautifulSoup(resp.text, "html.parser")

        article = soup.find("article") or soup.find("main") or soup
        for img in article.find_all("img"):
            src = img.get("src", "")
            if src and "favicon" not in src and "avatar" not in src and "logo" not in src and not src.startswith("data:"):
                images.append(src)

        if not images:
            og = soup.find("meta", property="og:image")
            if og and og.get("content"):
                images.append(og["content"])
    except Exception:
        pass
    return images


def download_image(url: str, images_dir: Path) -> tuple[Optional[str], Optional[bytes]]:
    """Download image to disk, return (relative_path, raw_bytes) or (None, None)."""
    try:
        resp = requests.get(url, timeout=15, headers={"User-Agent": "Mozilla/5.0"})
        resp.raise_for_status()
        content_type = resp.headers.get("content-type", "")
        if "image" in content_type:
            from PIL import Image as PILImage
            PILImage.open(BytesIO(resp.content)).verify()

            ext = ".png"
            if "jpeg" in content_type or "jpg" in content_type:
                ext = ".jpg"
            elif "gif" in content_type:
                ext = ".gif"
            elif "webp" in content_type:
                ext = ".webp"

            filename = hashlib.md5(url.encode()).hexdigest()[:12] + ext
            filepath = images_dir / filename
            filepath.write_bytes(resp.content)
            return f"images/{filename}", resp.content
    except Exception:
        pass
    return None, None


# ============================================================================
# Phase 2: AI Translation via GitHub Models API
# ============================================================================

TRANSLATION_SYSTEM_PROMPT = """You are a professional English-to-Japanese translator for a GitHub changelog slide deck.

Rules:
- Use dictionary form (辞書形), NOT polite form (です/ます)
- Keep English for product names: GitHub Copilot, GitHub Actions, CodeQL, Dependabot, etc.
- Keep English for technical terms when Japanese equivalent is uncommon: API, PR, CI/CD, OAuth, SSH, OIDC
- Do not translate version numbers, dates, or proper nouns
- Max ~50 characters per bullet for slide readability
- Each bullet must be a complete, standalone thought
- Write as a Japanese SE would for an internal slide deck

You will receive a JSON array of objects with "title" and "bullets" fields.
Return a JSON array of objects with "title_ja" and "bullets_ja" fields, in the same order.
Return ONLY valid JSON, no markdown fences or explanation."""


def translate_with_ai(entries: list[dict], token: str) -> bool:
    """Translate entries using GitHub Models API. Returns True on success."""
    # Prepare batches (max 10 entries per request to stay within token limits)
    batch_size = 10
    all_success = True

    for i in range(0, len(entries), batch_size):
        batch = entries[i:i + batch_size]
        payload = [{"title": e["title"], "bullets": e["bullets"]} for e in batch]

        try:
            resp = requests.post(
                MODELS_API_URL,
                headers={
                    "Authorization": f"Bearer {token}",
                    "Content-Type": "application/json",
                },
                json={
                    "model": MODELS_MODEL,
                    "messages": [
                        {"role": "system", "content": TRANSLATION_SYSTEM_PROMPT},
                        {"role": "user", "content": json.dumps(payload, ensure_ascii=False)},
                    ],
                    "temperature": 0.3,
                    "max_tokens": 4000,
                },
                timeout=60,
            )
            resp.raise_for_status()

            content = resp.json()["choices"][0]["message"]["content"]
            # Strip markdown fences if present
            content = content.strip()
            if content.startswith("```"):
                content = content.split("\n", 1)[1] if "\n" in content else content[3:]
                if content.endswith("```"):
                    content = content[:-3]
                content = content.strip()

            translations = json.loads(content)

            for j, trans in enumerate(translations):
                if i + j < len(entries):
                    entries[i + j]["title_ja"] = trans.get("title_ja", "")
                    entries[i + j]["bullets_ja"] = trans.get("bullets_ja", [])

            print(f"  Translated batch {i // batch_size + 1} ({len(batch)} entries)")

        except Exception as e:
            print(f"  Warning: AI translation failed for batch {i // batch_size + 1}: {e}")
            # Leave title_ja / bullets_ja empty — English-only output
            for j in range(len(batch)):
                if i + j < len(entries):
                    entries[i + j].setdefault("title_ja", "")
                    entries[i + j].setdefault("bullets_ja", [])
            all_success = False

    return all_success


# ============================================================================
# Phase 3: Markdown Generation
# ============================================================================


def generate_markdown(entries: list[dict], lang: str = "en", file_stem: str = "") -> str:
    today = datetime.now().strftime("%Y-%m-%d")

    # Jekyll frontmatter for GitHub Pages
    dates = [e["date"] for e in entries]
    date_from = min(dates).strftime("%b %d")
    date_to = max(dates).strftime("%b %d, %Y")
    if lang == "ja":
        fm_lines = [
            "---",
            f"title: 変更履歴 ({date_from}–{date_to})",
            "layout: default",
            "parent: News",
            "nav_exclude: true",
            "lang: ja",
        ]
        if file_stem:
            fm_lines.append(f"lang_pair: /10-news/11-changelog/11.01-english/{file_stem}-changelog.html")
        fm_lines.append("---\n")
        lines = ["\n".join(fm_lines)]
        lines += [f"# GitHub Changelog アップデート\n", f"**生成日**: {today}\n"]
    else:
        fm_lines = [
            "---",
            f"title: Changelog ({date_from}–{date_to})",
            "layout: default",
            "parent: News",
            "nav_order: 1",
            "lang: en",
        ]
        if file_stem:
            fm_lines.append(f"lang_pair: /10-news/11-changelog/11.02-japanese/{file_stem}-changelog.html")
        fm_lines.append("---\n")
        lines = ["\n".join(fm_lines)]
        lines += [f"# GitHub Changelog Updates\n", f"**Generated**: {today}\n"]

    for entry in entries:
        title = entry.get("title_ja", entry["title"]) if lang == "ja" else entry["title"]
        lines.append(f"## {title}\n")
        lines.append(f"- **Date**: {entry['date'].strftime('%Y-%m-%d')}")
        lines.append(f"- **Type**: {entry['type']}")
        lines.append(f"- **Category**: {entry['label']}")
        lines.append(f"- **Link**: {entry['link']}\n")

        bullets = entry.get("bullets_ja" if lang == "ja" else "bullets", [])
        for b in bullets:
            lines.append(f"- {b}")
        lines.append("")

    return "\n".join(lines)


# ============================================================================
# Phase 4: PPTX Generation
# ============================================================================


def _find_shapes(slide_xml):
    sp_tree = slide_xml.find(f".//{{{NS_P}}}spTree")
    return list(sp_tree) if sp_tree is not None else []


def _get_shape_name(shape_xml):
    cNvPr = shape_xml.find(f".//{{{NS_P}}}cNvPr")
    if cNvPr is None:
        cNvPr = shape_xml.find(f".//{{{NS_A}}}cNvPr")
    return cNvPr.get("name", "") if cNvPr is not None else ""


def _find_paragraphs(shape_xml):
    txBody = shape_xml.find(f".//{{{NS_P}}}txBody")
    if txBody is None:
        txBody = shape_xml.find(f".//{{{NS_A}}}txBody")
    return txBody.findall(f"{{{NS_A}}}p") if txBody is not None else []


def build_pptx(entries: list[dict], output_path: Path, lang: str = "ja"):
    """Build PPTX by cloning template slides and replacing text/images.

    Template slide map:
      0 = Title, 1 = TOC, 2 = Section header,
      3 = Content (IMPROVEMENT), 4 = Image, 5 = Content (RETIRED), 6 = Content (RELEASE)
    """
    prs = Presentation(str(TEMPLATE_PATH))

    tpl_xmls = [copy.deepcopy(s._element) for s in prs.slides]

    tpl_extra_rels = {}
    for idx, slide in enumerate(prs.slides):
        extra = {}
        for rel in slide.part.rels.values():
            if "slideLayout" not in rel.reltype and "notesSlide" not in rel.reltype:
                if rel.is_external:
                    extra[rel.rId] = (rel.reltype, rel.target_ref, True)
                else:
                    extra[rel.rId] = (rel.reltype, rel.target_part, False)
        if extra:
            tpl_extra_rels[idx] = extra

    tpl_pic_xml = None
    spTree = tpl_xmls[4].find(f".//{{{NS_P}}}spTree")
    for el in spTree:
        if el.tag == f"{{{NS_P}}}pic":
            tpl_pic_xml = copy.deepcopy(el)
            break

    def _clone_slide(prs, tpl_xml, tpl_idx):
        layout_idx = 0 if tpl_idx == 0 else 1
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        new_xml = copy.deepcopy(tpl_xml)
        el = slide._element
        for child in list(el):
            el.remove(child)
        for key, val in new_xml.attrib.items():
            el.set(key, val)
        for child in list(new_xml):
            el.append(child)
        if tpl_idx in tpl_extra_rels:
            for rId, (reltype, target, is_external) in tpl_extra_rels[tpl_idx].items():
                slide.part.rels._rels[rId] = _Relationship(
                    slide.part.rels._base_uri, rId, reltype,
                    target_mode="External" if is_external else "Internal",
                    target=target,
                )
        return slide

    def _set_title_text(slide_xml, text):
        for shape in _find_shapes(slide_xml):
            if "Title" in _get_shape_name(shape):
                paras = _find_paragraphs(shape)
                if paras:
                    runs = paras[0].findall(f"{{{NS_A}}}r")
                    if runs:
                        runs[0].find(f"{{{NS_A}}}t").text = text
                        for r in runs[1:]:
                            paras[0].remove(r)
                return

    def _set_subtitle_texts(slide_xml, texts):
        for shape in _find_shapes(slide_xml):
            if "Subtitle" in _get_shape_name(shape):
                paras = _find_paragraphs(shape)
                if paras:
                    runs = paras[0].findall(f"{{{NS_A}}}r")
                    for i, run in enumerate(runs):
                        t_elem = run.find(f"{{{NS_A}}}t")
                        if i < len(texts):
                            t_elem.text = texts[i]
                        else:
                            paras[0].remove(run)
                return

    def _set_content_bullets(slide_xml, bullets):
        for shape in _find_shapes(slide_xml):
            if "Content" in _get_shape_name(shape):
                txBody = shape.find(f".//{{{NS_P}}}txBody")
                if txBody is None:
                    continue
                paras = txBody.findall(f"{{{NS_A}}}p")
                if not paras:
                    continue
                bullet_tpl = copy.deepcopy(paras[0])
                sep_tpl = copy.deepcopy(paras[1]) if len(paras) > 1 else None
                for p in paras:
                    txBody.remove(p)
                for i, bullet_text in enumerate(bullets):
                    if i > 0 and sep_tpl is not None:
                        txBody.append(copy.deepcopy(sep_tpl))
                    p = copy.deepcopy(bullet_tpl)
                    runs = p.findall(f"{{{NS_A}}}r")
                    if runs:
                        runs[0].find(f"{{{NS_A}}}t").text = bullet_text
                        for r in runs[1:]:
                            p.remove(r)
                    txBody.append(p)
                return

    def _set_date_tile(slide_xml, date_obj):
        for shape in _find_shapes(slide_xml):
            if "Rounded" in _get_shape_name(shape):
                paras = _find_paragraphs(shape)
                if paras:
                    runs = paras[0].findall(f"{{{NS_A}}}r")
                    if runs and len(runs) >= 3:
                        runs[0].find(f"{{{NS_A}}}t").text = date_obj.strftime("%B").upper()
                        runs[1].find(f"{{{NS_A}}}t").text = " "
                        runs[2].find(f"{{{NS_A}}}t").text = date_obj.strftime("%d, %Y")
                        for r in runs[3:]:
                            paras[0].remove(r)
                    elif runs:
                        runs[0].find(f"{{{NS_A}}}t").text = date_obj.strftime("%B %d, %Y").upper()
                return

    def _set_category_tile(slide_xml, category_type):
        tpl_idx = CATEGORY_TPL_SLIDE.get(category_type, 3)
        src_tile = None
        rounded_count = 0
        for shape in _find_shapes(tpl_xmls[tpl_idx]):
            if "Rounded" in _get_shape_name(shape):
                rounded_count += 1
                if rounded_count == 2:
                    src_tile = copy.deepcopy(shape)
                    break
        if src_tile is None:
            return
        sp_tree = slide_xml.find(f".//{{{NS_P}}}spTree")
        rounded_count = 0
        for shape in list(_find_shapes(slide_xml)):
            if "Rounded" in _get_shape_name(shape):
                rounded_count += 1
                if rounded_count == 2:
                    sp_tree.remove(shape)
                    paras = src_tile.findall(f".//{{{NS_A}}}p")
                    if paras:
                        runs = paras[0].findall(f"{{{NS_A}}}r")
                        if runs:
                            runs[0].find(f"{{{NS_A}}}t").text = category_type.upper()
                            for r in runs[1:]:
                                paras[0].remove(r)
                    sp_tree.append(src_tile)
                    return

    def _remove_image(slide_xml):
        sp_tree = slide_xml.find(f".//{{{NS_P}}}spTree")
        if sp_tree is not None:
            for pic in sp_tree.findall(f"{{{NS_P}}}pic"):
                sp_tree.remove(pic)

    def _add_image_from_template(slide, pic_xml_tpl, image_bytes):
        from pptx.parts.image import ImagePart, Image as PptxImage
        try:
            pic_xml = copy.deepcopy(pic_xml_tpl)
            image = PptxImage.from_blob(image_bytes)
            image_part = ImagePart.new(prs.part.package, image)
            rId = slide.part.relate_to(image_part, RT.IMAGE)
            blip = pic_xml.find(f".//{{{NS_A}}}blip")
            if blip is not None:
                blip.set(f"{{{NS_R}}}embed", rId)
            spTree = slide._element.find(f".//{{{NS_P}}}spTree")
            spTree.append(pic_xml)
        except Exception:
            pass

    def _set_slide_notes(slide, text):
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.clear()
        tf.text = text

    # --- Remove all existing slides ---
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(f"{{{NS_R}}}id")
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    # --- Slide 0: Title ---
    title_xml = copy.deepcopy(tpl_xmls[0])
    dates = [e["date"] for e in entries]
    if dates:
        if lang == "ja":
            _set_subtitle_texts(title_xml, [
                min(dates).strftime("%-m月%-d日"), "〜 ",
                max(dates).strftime("%-m月%-d日"), " ", "のアップデート"
            ])
        else:
            _set_subtitle_texts(title_xml, [
                min(dates).strftime("%B"), " ",
                min(dates).strftime("%d"), " - ",
                max(dates).strftime("%B %d, %Y")
            ])
    _clone_slide(prs, title_xml, 0)

    # --- Slide 1: TOC ---
    groups = group_by_label(entries)
    toc_xml = copy.deepcopy(tpl_xmls[1])
    for shape in _find_shapes(toc_xml):
        if "Content" in _get_shape_name(shape):
            txBody = shape.find(f".//{{{NS_P}}}txBody")
            if txBody is None:
                continue
            paras = txBody.findall(f"{{{NS_A}}}p")
            toc_tpl = None
            empty_tpl = None
            for p in paras:
                runs = p.findall(f"{{{NS_A}}}r")
                if runs and toc_tpl is None:
                    toc_tpl = copy.deepcopy(p)
                elif not runs and empty_tpl is None:
                    empty_tpl = copy.deepcopy(p)
            for p in paras:
                txBody.remove(p)
            for i, (label, items) in enumerate(groups.items(), 1):
                if i > 1 and empty_tpl is not None:
                    txBody.append(copy.deepcopy(empty_tpl))
                p = copy.deepcopy(toc_tpl)
                runs = p.findall(f"{{{NS_A}}}r")
                if runs:
                    count_text = f"{len(items)}件" if lang == "ja" else str(len(items))
                    runs[0].find(f"{{{NS_A}}}t").text = f"{i} {label.title()} ("
                    if len(runs) > 1:
                        runs[1].find(f"{{{NS_A}}}t").text = count_text
                    if len(runs) > 2:
                        runs[2].find(f"{{{NS_A}}}t").text = ")"
                    for r in runs[3:]:
                        p.remove(r)
                txBody.append(p)
            break
    _clone_slide(prs, toc_xml, 1)

    # --- Per-category slides ---
    sections = [("Introduction", 0)]

    for label, items in groups.items():
        section_start = len(prs.slides)

        section_xml = copy.deepcopy(tpl_xmls[2])
        count_suffix = f"{len(items)}件" if lang == "ja" else str(len(items))
        _set_title_text(section_xml, f"{label.title()} ({count_suffix})")
        _clone_slide(prs, section_xml, 2)

        for entry in items:
            cat_type = entry.get("type", "Improvement")
            entry_date = entry["date"]
            title = entry.get("title_ja", entry["title"]) if lang == "ja" else entry["title"]
            bullets = entry.get("bullets_ja", []) if lang == "ja" else entry.get("bullets", [])

            content_xml = copy.deepcopy(tpl_xmls[3])
            _set_title_text(content_xml, title)
            _set_content_bullets(content_xml, bullets)
            _set_date_tile(content_xml, entry_date)
            _set_category_tile(content_xml, cat_type)
            summary_slide = _clone_slide(prs, content_xml, 3)

            notes_text = f"{entry['title']}\n{entry['link']}\n\n{entry.get('plain_text_en', '')}"
            _set_slide_notes(summary_slide, notes_text)

            image_bytes = entry.get("_image_bytes")
            if image_bytes and tpl_pic_xml is not None:
                img_xml = copy.deepcopy(tpl_xmls[4])
                _set_title_text(img_xml, title)
                _set_date_tile(img_xml, entry_date)
                _set_category_tile(img_xml, cat_type)
                _remove_image(img_xml)
                slide = _clone_slide(prs, img_xml, 4)
                _add_image_from_template(slide, tpl_pic_xml, image_bytes)

        sections.append((f"{label.title()} ({len(items)})", section_start))

    # --- Add PowerPoint sections ---
    prs_element = prs.part._element
    sld_id_lst = prs_element.find(f"{{{NS_P}}}sldIdLst")
    slide_ids = [el.get("id") for el in sld_id_lst]

    section_lst = None
    ext_lst = prs_element.find(f"{{{NS_P}}}extLst")
    if ext_lst is not None:
        for ext in ext_lst.findall(f"{{{NS_P}}}ext"):
            if ext.get("uri") == SECTION_EXT_URI:
                section_lst = ext.find(f"{{{NS_P14}}}sectionLst")
                if section_lst is None:
                    section_lst = etree.SubElement(ext, f"{{{NS_P14}}}sectionLst")
                else:
                    for child in list(section_lst):
                        section_lst.remove(child)
                break
        else:
            ext_el = etree.SubElement(ext_lst, f"{{{NS_P}}}ext")
            ext_el.set("uri", SECTION_EXT_URI)
            section_lst = etree.SubElement(ext_el, f"{{{NS_P14}}}sectionLst")

    if section_lst is None:
        ext_lst = etree.SubElement(prs_element, f"{{{NS_P}}}extLst")
        ext_el = etree.SubElement(ext_lst, f"{{{NS_P}}}ext")
        ext_el.set("uri", SECTION_EXT_URI)
        section_lst = etree.SubElement(ext_el, f"{{{NS_P14}}}sectionLst")

    for sec_name, sec_start in sections:
        sec_el = etree.SubElement(section_lst, f"{{{NS_P14}}}section")
        sec_el.set("name", sec_name)
        sec_el.set("id", "{" + str(uuid.uuid4()).upper() + "}")
        sld_id_lst_el = etree.SubElement(sec_el, f"{{{NS_P14}}}sldIdLst")
        sec_idx = sections.index((sec_name, sec_start))
        if sec_idx + 1 < len(sections):
            sec_end = sections[sec_idx + 1][1]
        else:
            sec_end = len(slide_ids)
        for sid in slide_ids[sec_start:sec_end]:
            sld_id_el = etree.SubElement(sld_id_lst_el, f"{{{NS_P14}}}sldId")
            sld_id_el.set("id", sid)

    prs.save(str(output_path))


# ============================================================================
# Main
# ============================================================================


def main():
    parser = argparse.ArgumentParser(
        description="Fetch GitHub changelog, translate with AI, and build markdown + PPTX"
    )
    parser.add_argument("--days", type=int, default=14, help="Days to look back (default: 14)")
    parser.add_argument("--output-dir", type=str, default=None, help="Output directory")
    args = parser.parse_args()

    output_dir = Path(args.output_dir) if args.output_dir else REPO_ROOT / "10-news" / "11-changelog"
    images_dir = output_dir / "images"
    en_dir = output_dir / "11.01-english"
    ja_dir = output_dir / "11.02-japanese"
    ppt_en_dir = output_dir / "11.03-ppt-english"
    ppt_ja_dir = output_dir / "11.04-ppt-japanese"

    for d in [images_dir, en_dir, ja_dir, ppt_en_dir, ppt_ja_dir]:
        d.mkdir(parents=True, exist_ok=True)

    # --- Phase 1: Fetch RSS ---
    print(f"Fetching changelog entries from the last {args.days} days...")
    entries = fetch_feed(args.days)
    print(f"Found {len(entries)} entries.")

    if not entries:
        print("No entries found. Exiting.")
        sys.exit(0)

    print("Extracting summaries and images...")
    for entry in entries:
        bullets, image_urls, plain_text = extract_summary_and_images(entry.pop("content_html", ""))
        entry["bullets"] = bullets
        entry["plain_text_en"] = plain_text

        if not image_urls:
            print(f"  Scraping images from: {entry['link'][:60]}...")
            image_urls = scrape_page_images(entry["link"])

        entry["_image_bytes"] = None
        for url in image_urls[:3]:
            rel_path, img_bytes = download_image(url, images_dir)
            if rel_path:
                entry["image_file"] = rel_path
                entry["_image_bytes"] = img_bytes
                break

    # --- Phase 2: AI Translation ---
    token = os.environ.get("GITHUB_TOKEN", "")
    if token:
        print("Translating to Japanese via GitHub Models API...")
        success = translate_with_ai(entries, token)
        if not success:
            print("Warning: Some translations failed. Japanese outputs may be incomplete.")
    else:
        print("Warning: GITHUB_TOKEN not set. Skipping AI translation (English-only output).")
        for entry in entries:
            entry["title_ja"] = ""
            entry["bullets_ja"] = []

    has_translations = any(e.get("title_ja") for e in entries)

    # --- Phase 3: Markdown ---
    dates = [e["date"] for e in entries]
    date_from = min(dates).strftime("%Y-%m-%d")
    date_to = max(dates).strftime("%Y-%m-%d")
    file_stem = f"{date_from}_to_{date_to}"

    en_md = generate_markdown(entries, "en", file_stem)
    en_path = en_dir / f"{file_stem}-changelog.md"
    en_path.write_text(en_md, encoding="utf-8")
    print(f"Saved English markdown: {en_path}")

    if has_translations:
        ja_md = generate_markdown(entries, "ja", file_stem)
        ja_path = ja_dir / f"{file_stem}-changelog.md"
        ja_path.write_text(ja_md, encoding="utf-8")
        print(f"Saved Japanese markdown: {ja_path}")

    # --- Phase 4: PPTX ---
    ppt_en_path = ppt_en_dir / f"{file_stem}-changelog.pptx"
    print("Building English PowerPoint...")
    build_pptx(entries, ppt_en_path, lang="en")
    print(f"Saved English PowerPoint: {ppt_en_path}")

    if has_translations:
        ppt_ja_path = ppt_ja_dir / f"{file_stem}-changelog.pptx"
        print("Building Japanese PowerPoint...")
        build_pptx(entries, ppt_ja_path, lang="ja")
        print(f"Saved Japanese PowerPoint: {ppt_ja_path}")
    else:
        print("Skipping Japanese outputs (no translations available).")

    print(f"\nDone! Generated {len(entries)} changelog entries.")


if __name__ == "__main__":
    main()
